# 다중 포맷 문서를 평문 텍스트로 변환하는 로더 (TRD §3)
import os
import re
from zipfile import BadZipFile

# DOCX/XLSX/PPTX는 모두 ZIP 컨테이너라, 구형 바이너리(.doc/.xls/.ppt)나 손상 파일이면
# 파서가 BadZipFile을 던진다. 영문 예외 대신 실제 감지 포맷을 담아 안내한다.
_ZIP_FORMAT_HINT = "유효한 {fmt} 파일이 아닙니다. {detected}"
_CONVERT_HINT = "최신 포맷이나 PDF로 변환해 업로드하세요."


def _sniff_format(path: str) -> str:
    # 파일 첫 바이트와 OLE 내부 스트림으로 실제 포맷을 추정해 한국어 한 문장으로 돌려준다.
    # 새 의존성 없이 기존 olefile만 사용한다. 진단/오류 메시지 보강용.
    try:
        with open(path, "rb") as f:
            head = f.read(512)
    except OSError:
        return "파일을 읽을 수 없습니다."
    if head[:4] == b"PK\x03\x04":
        return f"ZIP 컨테이너이지만 내부 구조가 손상된 것으로 보입니다. {_CONVERT_HINT}"
    if head[:4] == b"%PDF":
        return "실제로는 PDF 파일입니다. 확장자를 .pdf로 바꿔 업로드하세요."
    if head[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # OLE 복합 문서
        return f"구형 OLE 바이너리({_ole_subtype(path)})로 보입니다. {_CONVERT_HINT}"

    # 텍스트 계열 — 공공·기업 시스템이 보고서를 HTML/XML/CSV로 내보내며
    # 확장자만 .xls/.xlsx로 붙이는 경우가 흔하다. BOM은 제거 후 판별한다.
    sample = head.lstrip(b"\xef\xbb\xbf").lstrip()
    low = sample[:256].lower()
    if low.startswith(b"<?xml") and b"spreadsheet" in low:
        return "실제로는 Excel 2003 XML(SpreadsheetML) 파일입니다. 한글/오피스에서 .xlsx로 다시 저장해 업로드하세요."
    if low.startswith(b"<?xml"):
        return f"실제로는 XML 텍스트 파일입니다. {_CONVERT_HINT}"
    if low.startswith((b"<!doctype html", b"<html", b"<table")) or b"<table" in low:
        return "실제로는 HTML 표 파일입니다(시스템이 .xls로 내보낸 형식). 한글/오피스에서 .xlsx로 다시 저장하거나 CSV로 변환해 업로드하세요."
    text_msg = "실제로는 일반 텍스트/CSV로 보입니다. 확장자를 .txt 또는 .csv로 바꿔 업로드하세요."
    try:
        head.decode("utf-8")
        return text_msg
    except UnicodeDecodeError as ude:
        # 512바이트 경계에서 멀티바이트 문자가 잘렸을 뿐이면 여전히 텍스트로 본다.
        if ude.start >= len(head) - 4:
            return text_msg
    return f"알 수 없는 형식이거나 손상된 파일입니다(첫 바이트: {head[:8].hex(' ')}). {_CONVERT_HINT}"


def _ole_subtype(path: str) -> str:
    # OLE 복합 문서의 내부 스트림 이름으로 구형 오피스/HWP 종류를 구분한다.
    try:
        import olefile

        ole = olefile.OleFileIO(path)
        try:
            names = {e[0] if isinstance(e, list) else e for e in ole.listdir()}
            names |= {n.lstrip("\x05") for n in names}
        finally:
            ole.close()
    except Exception:
        return "구형 오피스/HWP"
    if {"Workbook", "Book"} & names:
        return "구형 Excel .xls"
    if "WordDocument" in names:
        return "구형 Word .doc"
    if "PowerPoint Document" in names:
        return "구형 PowerPoint .ppt"
    if {"FileHeader", "HwpSummaryInformation", "BodyText"} & names:
        return "한글 HWP"
    return "구형 오피스/HWP"


def load_document(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in (".txt", ".md", ".markdown"):
        return _read_text(path)
    if ext == ".pdf":
        return _load_pdf(path)
    if ext == ".docx":
        return _load_docx(path)
    if ext == ".pptx":
        return _load_pptx(path)
    if ext == ".xlsx":
        return _load_xlsx(path)
    if ext == ".hwp":
        return _load_hwp(path)
    if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        return _load_ocr(path)
    raise ValueError(f"지원하지 않는 포맷: {ext}")


def _read_text(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


# 텍스트 레이어가 이보다 적은 페이지는 이미지 기반(스캔)으로 보고 OCR로 보완한다.
_PDF_OCR_MIN_CHARS = 30


def _load_pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = [(p.extract_text() or "").strip() for p in reader.pages]
    # 스캔/이미지 PDF는 pypdf가 텍스트를 거의 못 뽑아 데이터가 비어버린다.
    # 텍스트가 빈약한 페이지만 렌더링해 OCR로 채운다(텍스트 페이지는 그대로).
    sparse = [i for i, t in enumerate(pages) if len(t) < _PDF_OCR_MIN_CHARS]
    if sparse:
        pages = _ocr_pdf_pages(path, pages, sparse)
    return "\n".join(pages)


def _ocr_pdf_pages(path, pages, sparse):
    # sparse 페이지를 pypdfium2로 300DPI 렌더링해 OCR한다(300DPI+PSM6가 정확도 최적).
    # 렌더러/Tesseract 미설치 시 조용히 빈약한 데이터를 내지 않고 명확한 안내(ValueError)로 폴백.
    try:
        import pypdfium2 as pdfium
    except ModuleNotFoundError as e:
        raise ValueError(
            "이미지 기반(스캔) PDF로 보입니다. PDF 렌더러(pypdfium2)가 설치되지 않아 "
            "OCR할 수 없습니다. 오프라인 반입하거나 텍스트 PDF로 변환해 업로드하세요."
        ) from e
    pdf = pdfium.PdfDocument(path)
    try:
        for i in sparse:
            img = pdf[i].render(scale=300 / 72).to_pil()
            pages[i] = _ocr_image(img).strip()
    finally:
        pdf.close()
    return pages


def _load_docx(path):
    import docx
    from docx.opc.exceptions import PackageNotFoundError

    # 비-ZIP/구형 .doc는 BadZipFile 또는 PackageNotFoundError로 나타난다.
    try:
        d = docx.Document(path)
    except (BadZipFile, PackageNotFoundError) as e:
        raise ValueError(_ZIP_FORMAT_HINT.format(fmt="DOCX", detected=_sniff_format(path))) from e
    return "\n".join(p.text for p in d.paragraphs)


def _load_pptx(path):
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(path)
    except (BadZipFile, PackageNotFoundError) as e:
        raise ValueError(_ZIP_FORMAT_HINT.format(fmt="PPTX", detected=_sniff_format(path))) from e
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def _load_xlsx(path):
    from openpyxl import load_workbook

    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except BadZipFile as e:
        raise ValueError(_ZIP_FORMAT_HINT.format(fmt="XLSX", detected=_sniff_format(path))) from e
    out = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out)


def _load_hwp(path):
    # HWP는 경량 best-effort: olefile로 미리보기 텍스트(PrvText) 스트림만 추출한다.
    # 본문 완전 추출은 미지원(망분리·경량 정책)이며, 추출 실패 시 빈 결과 대신
    # 한계를 명확히 알려 사용자가 대안(PDF/DOCX 변환)을 택하게 한다.
    import olefile

    if not olefile.isOleFile(path):
        raise ValueError("유효한 HWP(OLE) 파일이 아닙니다.")
    try:
        ole = olefile.OleFileIO(path)
    except Exception as e:
        # OLE 매직은 맞지만 내부 구조가 깨진 파일 → 깔끔한 400으로 안내(500 방지).
        raise ValueError("손상된 HWP(OLE) 파일입니다. PDF 또는 DOCX로 변환해 업로드하세요.") from e
    try:
        if ole.exists("PrvText"):
            text = ole.openstream("PrvText").read().decode("utf-16-le", errors="ignore").strip()
            if text:
                return text
    finally:
        ole.close()
    raise ValueError(
        "HWP 본문을 추출하지 못했습니다(미리보기 텍스트 없음). "
        "현재 경량 로더는 PrvText 미리보기만 지원합니다. 본문이 필요하면 "
        "문서를 PDF 또는 DOCX로 변환해 업로드하세요."
    )


# Tesseract가 한글을 글자 단위로 떼어내며 넣는 공백('건 설 산 업')을 붙인다.
# 한 글자+공백이 3개 이상 연속되는 구간만 제거하므로, 정상 단어 사이 공백
# ('검토 보고')이나 2글자 이상 토큰('미지 급')은 건드리지 않는다(오병합 방지).
_OCR_CHAR_RUN = re.compile(r"(?:[가-힣] ){2,}[가-힣]")


def _collapse_ocr_spacing(text: str) -> str:
    return _OCR_CHAR_RUN.sub(lambda m: m.group().replace(" ", ""), text)


def _ocr_image(img):
    # PIL 이미지를 Tesseract(한국어+영어)로 OCR한다. 이미지 파일 로더와 PDF 페이지
    # OCR 폴백이 공유한다. 바이너리 경로는 PATH 또는 TESSERACT_CMD, 언어데이터는
    # TESSDATA_PREFIX로 지정한다(머신별 경로 하드코딩 회피). 의존성·바이너리 부재는 명확히 안내.
    try:
        import pytesseract
    except ModuleNotFoundError as e:
        raise ValueError(
            "OCR 의존성(pytesseract)이 설치되지 않았습니다. 오프라인 설치 패키지로 반입하세요."
        ) from e

    cmd = os.environ.get("TESSERACT_CMD")
    if cmd:
        pytesseract.pytesseract.tesseract_cmd = cmd

    try:
        # PSM 6(균일 텍스트 블록)이 기본 PSM 3(자동 분할)보다 한국어 문서 페이지에서
        # 훨씬 정확하다(실측 19%→92%). 문서/표 위주 입력에 적합.
        raw = pytesseract.image_to_string(img, lang="kor+eng", config="--psm 6")
    except pytesseract.TesseractNotFoundError as e:
        raise ValueError(
            "Tesseract 실행파일을 찾을 수 없습니다. TESSERACT_CMD 환경변수로 "
            "경로를 지정하거나 오프라인 설치하세요."
        ) from e
    return _collapse_ocr_spacing(raw)


def _load_ocr(path):
    # 이미지 파일에서 텍스트 추출(Tesseract). PDF 스캔 페이지도 같은 _ocr_image를 쓴다.
    try:
        from PIL import Image
    except ModuleNotFoundError as e:
        raise ValueError(
            "OCR 의존성(Pillow)이 설치되지 않았습니다. 오프라인 설치 패키지로 반입하세요."
        ) from e
    return _ocr_image(Image.open(path))
